#!/usr/bin/env python3
"""
Genera presentación PPT premium: Historia de la Inteligencia Artificial
Desde los precursores hasta las grandes empresas actuales con CEOs y biografías
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ─── COLOR PALETTE (Premium Dark Theme) ───
BG_DARK = RGBColor(0x0D, 0x0D, 0x1A)
BG_SLIDE = RGBColor(0x13, 0x13, 0x2D)
ACCENT_CYAN = RGBColor(0x00, 0xE5, 0xFF)
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
ACCENT_GOLD = RGBColor(0xFF, 0xD7, 0x00)
ACCENT_NEON = RGBColor(0x39, 0xFF, 0x14)
ACCENT_CORAL = RGBColor(0xFF, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x33, 0x33, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=BG_DARK):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_banner(slide, left=0, top=0, width=None, height=Inches(0.08), color=ACCENT_CYAN):
    """Add a decorative banner line."""
    if width is None:
        width = prs.slide_width
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_gradient_banner(slide, left=Inches(0.5), top=Inches(0.3), width=None, height=Inches(0.05)):
    """Add a thin gradient-like banner (simulated with multiple colored segments)."""
    if width is None:
        width = prs.slide_width - Inches(1)
    segments = 4
    colors = [ACCENT_CYAN, ACCENT_GOLD, ACCENT_CORAL, ACCENT_PURPLE]
    seg_w = width // segments
    for i in range(segments):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + seg_w * i, top, seg_w, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i]
        shape.line.fill.background()

def add_title_text(slide, text, left=Inches(1), top=Inches(0.5), width=Inches(11), height=Inches(1.2), font_size=Pt(44), color=WHITE, bold=True, alignment=PP_ALIGN.LEFT):
    """Add a title text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Segoe UI'
    p.alignment = alignment
    return txBox

def add_subtitle_text(slide, text, left=Inches(1), top=Inches(1.8), width=Inches(11), height=Inches(0.8), font_size=Pt(22), color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a subtitle text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Segoe UI Light'
    p.alignment = alignment
    return txBox

def add_body_text(slide, paragraphs_list, left=Inches(1), top=Inches(2.5), width=Inches(11), height=Inches(4.5), font_size=Pt(18), color=LIGHT_GRAY, line_spacing=Pt(28)):
    """Add multi-paragraph body text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, para_text in enumerate(paragraphs_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para_text
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.space_after = Pt(8)
        p.line_spacing = line_spacing
    return txBox

def add_bullet_slide(slide, title, bullets, title_color=ACCENT_CYAN, accent_bullet_char="▸"):
    """Standard bullet-point slide."""
    add_bg(slide, BG_DARK)
    add_gradient_banner(slide)
    add_title_text(slide, title, top=Inches(0.5), color=title_color, font_size=Pt(40))
    body = add_body_text(slide, [f"{accent_bullet_char} {b}" for b in bullets], top=Inches(1.8), font_size=Pt(20))
    # Make first few words of each bullet cyan
    for para in body.text_frame.paragraphs:
        for run in para.runs:
            if len(run.text) > 2:
                # find colon and color differently
                pass
    return slide

def add_side_banner(slide, left=Inches(0.3), top=Inches(2), width=Inches(0.06), height=Inches(3.5), color=ACCENT_CYAN):
    """Vertical accent line on the left."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_page_number(slide, num, total):
    """Add page number at bottom right."""
    txBox = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = MED_GRAY
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.RIGHT

TOTAL_SLIDES = 28  # Will be updated

# ============================================================
# SLIDE 1: TITLE SLIDE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BG_DARK)

# Top banner
add_banner(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), ACCENT_CYAN)
add_banner(slide, Inches(0), Inches(0.12), Inches(13.333), Inches(0.04), ACCENT_GOLD)
add_banner(slide, Inches(0), Inches(0.16), Inches(13.333), Inches(0.04), ACCENT_CORAL)
add_banner(slide, Inches(0), Inches(0.20), Inches(13.333), Inches(0.04), ACCENT_PURPLE)

# Main title
add_title_text(slide, "HISTORIA DE LA", top=Inches(1.5), font_size=Pt(56), color=WHITE, alignment=PP_ALIGN.CENTER, width=Inches(12))
add_title_text(slide, "INTELIGENCIA ARTIFICIAL", top=Inches(2.3), font_size=Pt(56), color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER, width=Inches(12))

# Subtitle
add_subtitle_text(slide, "De los precursores a las grandes empresas del siglo XXI", top=Inches(3.3), font_size=Pt(26), color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, width=Inches(12), height=Inches(0.6))

# Decorative line
add_banner(slide, Inches(4), Inches(4.1), Inches(5.333), Inches(0.03), ACCENT_GOLD)

# Bottom info
add_subtitle_text(slide, "Presentación Premium · 2026", top=Inches(4.5), font_size=Pt(18), color=MED_GRAY, alignment=PP_ALIGN.CENTER, width=Inches(12))
add_subtitle_text(slide, "Generado por JARVIS · Agente NEO", top=Inches(6.8), font_size=Pt(14), color=DARK_GRAY, alignment=PP_ALIGN.CENTER, width=Inches(12))

# Bottom banner
add_banner(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.04), ACCENT_PURPLE)
add_banner(slide, Inches(0), Inches(7.39), Inches(13.333), Inches(0.04), ACCENT_CORAL)
add_banner(slide, Inches(0), Inches(7.43), Inches(13.333), Inches(0.04), ACCENT_GOLD)
add_banner(slide, Inches(0), Inches(7.47), Inches(13.333), Inches(0.08), ACCENT_CYAN)

add_page_number(slide, 1, TOTAL_SLIDES)

# ============================================================
# SLIDE 2: ÍNDICE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_title_text(slide, "ÍNDICE", top=Inches(0.5), color=ACCENT_CYAN, font_size=Pt(44))

indice_items = [
    ("01", "¿Qué es la Inteligencia Artificial?", ACCENT_CYAN),
    ("02", "Los Precursores: Turing, McCarthy, Minsky", ACCENT_GOLD),
    ("03", "La Conferencia de Dartmouth 1956", ACCENT_PURPLE),
    ("04", "Primeros Pasos (1950s-1960s)", ACCENT_CORAL),
    ("05", "El Primer Invierno de la IA (1970s)", ACCENT_NEON),
    ("06", "Sistemas Expertos y Resurgimiento (1980s)", ACCENT_CYAN),
    ("07", "Deep Blue y los 90s", ACCENT_GOLD),
    ("08", "Big Data y Machine Learning (2000s)", ACCENT_PURPLE),
    ("09", "La Revolución del Deep Learning (2010s)", ACCENT_CORAL),
    ("10", "Grandes Empresas de IA y sus CEOs", ACCENT_NEON),
    ("11", "El Futuro de la IA", ACCENT_CYAN),
]

y_start = Inches(1.8)
for i, (num, title, color) in enumerate(indice_items):
    y = y_start + Inches(i * 0.48)
    # Number
    txBox = slide.shapes.add_textbox(Inches(1.5), y, Inches(0.8), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Segoe UI'
    # Title
    txBox = slide.shapes.add_textbox(Inches(2.5), y, Inches(9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = 'Segoe UI Light'

add_page_number(slide, 2, TOTAL_SLIDES)

# ============================================================
# SLIDE 3: ¿QUÉ ES LA IA?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

bullets = [
    "Rama de la informática que busca crear sistemas capaces de realizar tareas que normalmente requieren inteligencia humana",
    "Incluye capacidades como: aprendizaje, razonamiento, percepción, comprensión del lenguaje natural y creatividad",
    "Se divide en IA Débil (específica) e IA Fuerte (general): la primera realiza tareas concretas, la segunda igualaría la inteligencia humana",
    "Hoy vivimos en la era de la IA Débil especializada: asistentes virtuales, coches autónomos, diagnóstico médico, generación de contenido",
    "El objetivo final: una Inteligencia Artificial General (AGI) que pueda razonar, aprender y adaptarse a cualquier tarea intelectual humana",
]
add_bullet_slide(slide, "¿QUÉ ES LA INTELIGENCIA ARTIFICIAL?", bullets)
add_page_number(slide, 3, TOTAL_SLIDES)

# ============================================================
# SLIDE 4: PRECURSORES - ALAN TURING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_side_banner(slide, color=ACCENT_GOLD)
add_title_text(slide, "ALAN TURING (1912-1954)", color=ACCENT_GOLD, font_size=Pt(40))
add_subtitle_text(slide, "El padre de la computación y la inteligencia artificial", top=Inches(1.4), font_size=Pt(22), color=ACCENT_CYAN)

bullets = [
    "Matemático, criptógrafo y filósofo británico. Descifró el código Enigma en la Segunda Guerra Mundial.",
    "En 1936 publicó \"On Computable Numbers\", donde introdujo el concepto de la Máquina de Turing, fundamento de la computación moderna.",
    "En 1950 publicó \"Computing Machinery and Intelligence\", donde formuló la famosa pregunta: ¿Pueden las máquinas pensar?",
    "Propuso el Test de Turing: si una máquina puede conversar sin ser distinguida de un humano, demuestra inteligencia.",
    "Su trabajo sentó las bases teóricas de la IA décadas antes de que el término existiera. Un visionario adelantado a su tiempo.",
]
add_body_text(slide, bullets, top=Inches(2.5), font_size=Pt(18))

# Quote box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.5), Inches(4.3), Inches(1.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
shape.line.color.rgb = ACCENT_GOLD
shape.line.width = Pt(1)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"Solo podemos ver poco del futuro, pero lo suficiente para saber que hay mucho por hacer."'
p.font.size = Pt(13)
p.font.color.rgb = ACCENT_GOLD
p.font.italic = True
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = '— Alan Turing'
p2.font.size = Pt(11)
p2.font.color.rgb = MED_GRAY
p2.font.name = 'Segoe UI'
p2.alignment = PP_ALIGN.RIGHT

add_page_number(slide, 4, TOTAL_SLIDES)

# ============================================================
# SLIDE 5: JOHN MCCARTHY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_side_banner(slide, color=ACCENT_PURPLE)
add_title_text(slide, "JOHN McCARTHY (1927-2011)", color=ACCENT_PURPLE, font_size=Pt(40))
add_subtitle_text(slide, "El hombre que acuñó el término 'Inteligencia Artificial'", top=Inches(1.4), font_size=Pt(22), color=ACCENT_CYAN)

bullets = [
    "Informático y científico cognitivo estadounidense. Ganador del Premio Turing en 1971.",
    "En 1956 acuñó el término \"Inteligencia Artificial\" para la Conferencia de Dartmouth, que organizó junto a Minsky, Shannon y Rochester.",
    "Inventó el lenguaje de programación LISP en 1958, el lenguaje estándar de la IA durante décadas y base del procesamiento simbólico.",
    "Pionero en el concepto de \"time-sharing\" (tiempo compartido) en computación, permitiendo que múltiples usuarios usaran una misma máquina.",
    "Fundó el Stanford AI Laboratory (SAIL) en 1963 y el MIT AI Lab, dos de los centros de investigación en IA más importantes del mundo.",
    "Propuso el \"Advice Taker\", uno de los primeros diseños de un sistema de IA completo, que influenció toda la IA simbólica posterior.",
]
add_body_text(slide, bullets, top=Inches(2.3), font_size=Pt(17))

# Quote
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.5), Inches(4.3), Inches(1.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
shape.line.color.rgb = ACCENT_PURPLE
shape.line.width = Pt(1)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"El estudio de la IA es hacer que las máquinas se comporten de formas que serían llamadas inteligentes si un humano lo hiciera."'
p.font.size = Pt(12)
p.font.color.rgb = ACCENT_PURPLE
p.font.italic = True
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = '— John McCarthy'
p2.font.size = Pt(11)
p2.font.color.rgb = MED_GRAY
p2.font.name = 'Segoe UI'
p2.alignment = PP_ALIGN.RIGHT

add_page_number(slide, 5, TOTAL_SLIDES)

# ============================================================
# SLIDE 6: MARVIN MINSKY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_side_banner(slide, color=ACCENT_CORAL)
add_title_text(slide, "MARVIN MINSKY (1927-2016)", color=ACCENT_CORAL, font_size=Pt(40))
add_subtitle_text(slide, "Pionero de las redes neuronales y la ciencia cognitiva", top=Inches(1.4), font_size=Pt(22), color=ACCENT_CYAN)

bullets = [
    "Científico cognitivo, informático y filósofo estadounidense. Cofundador del MIT AI Lab y del MIT Media Lab.",
    "En 1951 construyó SNARC, la primera máquina de aprendizaje basada en redes neuronales artificiales (con 40 sinapsis simuladas).",
    "Co-organizó la Conferencia de Dartmouth de 1956 y fue una figura central en el desarrollo temprano de la IA.",
    "Autor de \"Perceptrons\" (1969) con Seymour Papert, que analizó las limitaciones de las redes neuronales de una capa, influyendo en el primer invierno de la IA.",
    "Ganador del Premio Turing en 1969. También escribió \"The Society of Mind\" (1986), proponiendo que la mente es una sociedad de agentes simples.",
    "Trabajó en robótica, visión artificial, procesamiento del lenguaje natural y representación del conocimiento. Un verdadero polímata de la IA.",
]
add_body_text(slide, bullets, top=Inches(2.3), font_size=Pt(17))

# Quote
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.5), Inches(4.3), Inches(1.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
shape.line.color.rgb = ACCENT_CORAL
shape.line.width = Pt(1)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"La pregunta no es si las máquinas piensan, sino si los hombres lo hacen."'
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_CORAL
p.font.italic = True
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = '— Marvin Minsky'
p2.font.size = Pt(11)
p2.font.color.rgb = MED_GRAY
p2.font.name = 'Segoe UI'
p2.alignment = PP_ALIGN.RIGHT

add_page_number(slide, 6, TOTAL_SLIDES)

# ============================================================
# SLIDE 7: OTROS PIONEROS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_title_text(slide, "OTROS PRECURSORES FUNDAMENTALES", color=ACCENT_NEON, font_size=Pt(38))

pioneers = [
    ("Claude Shannon (1916-2001)", ACCENT_GOLD, "Padre de la teoría de la información. En 1950 creó \"Theseus\", un ratón mecánico que aprendía a resolver laberintos. Su trabajo matemático sobre la información digital es la base de toda la computación moderna."),
    ("Herbert Simon (1916-2001)", ACCENT_CYAN, "Economista y politólogo, ganador del Nobel de Economía (1978). Junto a Allen Newell creó Logic Theorist (1955), considerado el primer programa de IA. Acuñó el concepto de 'racionalidad limitada'."),
    ("Allen Newell (1927-1992)", ACCENT_PURPLE, "Junto a Simon desarrolló Logic Theorist y GPS (General Problem Solver). Creó SOAR, una arquitectura cognitiva unificada. Ganador del Premio Turing (1975). Fundamental en IA simbólica."),
    ("Frank Rosenblatt (1928-1971)", ACCENT_CORAL, "Creó el Perceptrón en 1957, la primera red neuronal artificial implementada en hardware (Mark I Perceptron). Su trabajo fue redescubierto décadas después como base del deep learning moderno."),
    ("Norbert Wiener (1894-1964)", ACCENT_GOLD, "Fundador de la cibernética, el estudio del control y comunicación en máquinas y seres vivos. Su libro 'Cybernetics' (1948) influyó profundamente en la IA, la robótica y la automatización."),
]

y = Inches(1.8)
for name, color, desc in pioneers:
    # Name
    txBox = slide.shapes.add_textbox(Inches(1), y, Inches(11.5), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(20)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Segoe UI'
    # Description
    txBox = slide.shapes.add_textbox(Inches(1.5), y + Inches(0.35), Inches(10.5), Inches(0.55))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = 'Segoe UI Light'
    y += Inches(1.05)

add_page_number(slide, 7, TOTAL_SLIDES)

# ============================================================
# SLIDE 8: CONFERENCIA DE DARTMOUTH 1956
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_title_text(slide, "CONFERENCIA DE DARTMOUTH 1956", color=ACCENT_GOLD, font_size=Pt(40))
add_subtitle_text(slide, "El nacimiento oficial de la Inteligencia Artificial como campo científico", top=Inches(1.3), font_size=Pt(20), color=ACCENT_CYAN)

bullets = [
    "Celebrada en el Dartmouth College (New Hampshire, EE.UU.) durante el verano de 1956, organizada por John McCarthy.",
    "Participantes: John McCarthy, Marvin Minsky, Claude Shannon, Nathaniel Rochester, Herbert Simon, Allen Newell, entre otros.",
    "Propusieron: 'Todo aspecto del aprendizaje o cualquier otra característica de la inteligencia puede ser descrito con tal precisión que puede construirse una máquina para simularlo'.",
    "Aquí se acuñó oficialmente el término 'Inteligencia Artificial' y se establecieron las bases del campo para las siguientes décadas.",
    "Se propusieron temas clave: redes neuronales, creatividad computacional, procesamiento del lenguaje natural y razonamiento abstracto.",
    "Aunque los resultados concretos fueron modestos, la conferencia creó la comunidad científica que impulsaría la IA durante 70 años.",
]
add_body_text(slide, bullets, top=Inches(2.3), font_size=Pt(18))

# Highlight box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(6.0), Inches(10.3), Inches(1.0))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
shape.line.color.rgb = ACCENT_GOLD
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🎯  1956: Año oficial del nacimiento de la IA. La conferencia duró 8 semanas y contó con solo 10 participantes, pero cambió el mundo para siempre."
p.font.size = Pt(18)
p.font.color.rgb = ACCENT_GOLD
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

add_page_number(slide, 8, TOTAL_SLIDES)

# ============================================================
# SLIDE 9: PRIMEROS PASOS (1950s-1960s)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_title_text(slide, "1950s-1960s: PRIMEROS PASOS", color=ACCENT_CYAN, font_size=Pt(40))

milestones = [
    ("1951", "SNARC", ACCENT_CYAN, "Marvin Minsky construye la primera red neuronal artificial con hardware. Simulaba el comportamiento de una rata en un laberinto usando 40 'sinapsis'."),
    ("1955", "Logic Theorist", ACCENT_GOLD, "Newell y Simon crean el primer programa de IA. Demostró 38 de los primeros 52 teoremas de Principia Mathematica, incluso encontrando una prueba más elegante."),
    ("1958", "LISP", ACCENT_PURPLE, "John McCarthy inventa LISP, el lenguaje de programación estándar de la IA durante décadas. Introdujo árboles, recolección de basura y procesamiento simbólico."),
    ("1961", "Unimate", ACCENT_CORAL, "Primer robot industrial instalado en General Motors. Marcó el inicio de la robótica industrial, manejando piezas de metal caliente."),
    ("1966", "ELIZA", ACCENT_NEON, "Joseph Weizenbaum crea ELIZA en el MIT, el primer chatbot. Simulaba un terapeuta rogeriano y sorprendió por su capacidad de mantener conversaciones."),
    ("1969", "Shakey", ACCENT_CYAN, "SRI International desarrolla Shakey, el primer robot móvil con capacidad de razonamiento. Podía navegar, planificar rutas y manipular objetos simples."),
]

y = Inches(1.7)
for year, name, color, desc in milestones:
    # Year box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(1.0), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = year
    p.font.size = Pt(16)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER
    # Name
    txBox = slide.shapes.add_textbox(Inches(2.0), y, Inches(2.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(17)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Segoe UI'
    # Description
    txBox = slide.shapes.add_textbox(Inches(4.6), y, Inches(8), Inches(0.67))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = 'Segoe UI Light'
    y += Inches(0.85)

add_page_number(slide, 9, TOTAL_SLIDES)

# ============================================================
# SLIDE 10: PRIMER INVIERNO DE LA IA (1970s)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_title_text(slide, "1970s: EL PRIMER INVIERNO DE LA IA", color=ACCENT_CORAL, font_size=Pt(40))
add_subtitle_text(slide, "Cuando las expectativas superaron a la realidad", top=Inches(1.3), font_size=Pt(20), color=LIGHT_GRAY)

bullets = [
    "Tras el optimismo inicial de los 60s, los avances en IA no cumplieron las expectativas desmesuradas generadas por los investigadores.",
    "El informe Lighthill (1973) en Reino Unido concluyó que la IA no lograría sus objetivos 'grandiosos', provocando recortes masivos de financiación gubernamental.",
    "DARPA (agencia militar de EE.UU.) también redujo drásticamente sus inversiones en investigación de IA básica.",
    "El libro 'Perceptrons' de Minsky y Papert (1969) demostró limitaciones de las redes neuronales de una capa, frenando la investigación en conexionismo.",
    "A pesar del invierno, surgieron avances importantes: PROLOG (1972) en Francia, el sistema experto MYCIN (1976) para diagnóstico médico.",
    "Lección clave: la IA sufre ciclos de hype y decepción. El invierno purgó promesas exageradas pero sentó bases para enfoques más realistas.",
]
add_body_text(slide, bullets, top=Inches(2.3), font_size=Pt(18))

add_page_number(slide, 10, TOTAL_SLIDES)

# ============================================================
# SLIDE 11: SISTEMAS EXPERTOS (1980s)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_title_text(slide, "1980s: SISTEMAS EXPERTOS Y RESURGIMIENTO", color=ACCENT_NEON, font_size=Pt(38))

bullets = [
    "Los Sistemas Expertos basados en reglas (if-then) se convirtieron en la primera aplicación comercialmente exitosa de la IA.",
    "XCON/R1 (1980): desarrollado por DEC para configurar pedidos de computadoras. Ahorró $40 millones anuales y demostró el ROI de la IA.",
    "MYCIN y DENDRAL: sistemas expertos para diagnóstico médico y análisis químico. Superaban a expertos humanos en sus dominios específicos.",
    "Japón lanza el proyecto 'Fifth Generation' (1982): ambiciosa iniciativa de 10 años para crear computadoras con IA, inspirada en PROLOG.",
    "Resurge el interés en redes neuronales: John Hopfield (1982) publica su modelo de red neuronal recurrente, revitalizando el conexionismo.",
    "El algoritmo de backpropagation se populariza (Rumelhart, Hinton, Williams, 1986), permitiendo entrenar redes neuronales multicapa eficientemente.",
    "Segundo 'boom' de la IA: inversión corporativa masiva en sistemas expertos. Para 1988, el mercado de IA alcanzó miles de millones de dólares.",
]
add_body_text(slide, bullets, top=Inches(1.8), font_size=Pt(17))

add_page_number(slide, 11, TOTAL_SLIDES)

# ============================================================
# SLIDE 12: DEEP BLUE Y LOS 90s
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_title_text(slide, "1990s: DEEP BLUE Y LA IA PRÁCTICA", color=ACCENT_GOLD, font_size=Pt(40))

bullets = [
    "1997: Deep Blue de IBM derrota al campeón mundial de ajedrez Garry Kasparov. Hito simbólico que demostró que las máquinas podían vencer a los mejores humanos en tareas intelectuales.",
    "Segundo invierno de la IA (finales 80s-principios 90s): los sistemas expertos resultaron difíciles de mantener y escalar. Nueva ola de recortes.",
    "La IA se vuelve 'invisible': algoritmos de IA se integran en productos cotidianos (motores de búsqueda, sistemas de recomendación, logística).",
    "1995: Richard Wallace crea A.L.I.C.E., chatbot ganador del Premio Loebner, basado en pattern matching con AIML.",
    "1998: Yann LeCun publica LeNet-5, red neuronal convolucional para reconocimiento de dígitos manuscritos. Precursor del deep learning moderno.",
    "1999: Sony lanza AIBO, el perro robot con IA. Primera mascota robótica comercial con capacidades de aprendizaje y personalidad.",
]
add_body_text(slide, bullets, top=Inches(2.0), font_size=Pt(18))

add_page_number(slide, 12, TOTAL_SLIDES)

# ============================================================
# SLIDE 13: BIG DATA Y MACHINE LEARNING (2000s)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_title_text(slide, "2000s: BIG DATA Y MACHINE LEARNING", color=ACCENT_PURPLE, font_size=Pt(38))

bullets = [
    "La explosión de internet generó cantidades masivas de datos ('Big Data'), el combustible perfecto para los algoritmos de aprendizaje automático.",
    "Google se convierte en la empresa líder en IA aplicada: PageRank, Google Translate (2006), sistemas de recomendación y búsqueda semántica.",
    "2006: Geoffrey Hinton acuña el término 'Deep Learning' y publica un artículo fundamental sobre entrenamiento de redes neuronales profundas usando pre-entrenamiento capa por capa.",
    "2007: Fei-Fei Li inicia ImageNet, base de datos de 14 millones de imágenes etiquetadas que impulsaría la revolución del computer vision.",
    "Amazon Web Services (AWS) democratiza el acceso a computación masiva, permitiendo a startups entrenar modelos de ML sin infraestructura propia.",
    "2009: comienza la competición anual ImageNet (ILSVRC), que catalizaría avances en visión artificial durante la siguiente década.",
    "El Machine Learning pasa de ser una curiosidad académica a una herramienta empresarial esencial: detección de fraude, recomendaciones, predicción de demanda.",
]
add_body_text(slide, bullets, top=Inches(1.8), font_size=Pt(17))

add_page_number(slide, 13, TOTAL_SLIDES)

# ============================================================
# SLIDE 14: REVOLUCIÓN DEL DEEP LEARNING (2010s)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_title_text(slide, "2010s: LA REVOLUCIÓN DEL DEEP LEARNING", color=ACCENT_CYAN, font_size=Pt(38))

bullets = [
    "2012: AlexNet (Alex Krizhevsky, Ilya Sutskever, Geoffrey Hinton) gana ImageNet por un margen histórico usando GPUs. El momento 'Big Bang' del deep learning.",
    "Las GPUs de NVIDIA se convierten en el hardware estándar para entrenar redes neuronales profundas, acelerando el entrenamiento 100x.",
    "2014: Ian Goodfellow inventa las GANs (Generative Adversarial Networks). Dos redes compitiendo entre sí revolucionan la generación de imágenes realistas.",
    "2014: DeepMind (adquirida por Google) presenta redes neuronales que aprenden a jugar videojuegos de Atari desde cero, solo viendo los píxeles.",
    "2016: AlphaGo de DeepMind derrota a Lee Sedol, campeón mundial de Go. Un hito considerado 10 años antes de lo previsto por los expertos.",
    "2017: Google presenta Transformers ('Attention Is All You Need'), la arquitectura que cambiaría para siempre el NLP y habilitaría los LLMs modernos.",
    "2018: OpenAI lanza GPT-1, BERT de Google revoluciona el NLP. Comienza la era de los modelos de lenguaje pre-entrenados a gran escala.",
]
add_body_text(slide, bullets, top=Inches(1.8), font_size=Pt(17))

add_page_number(slide, 14, TOTAL_SLIDES)

# ============================================================
# SLIDE 15: EMPRESAS LÍDERES - OPENAI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_side_banner(slide, color=ACCENT_NEON)
add_title_text(slide, "OPENAI", color=ACCENT_NEON, font_size=Pt(44))
add_subtitle_text(slide, "Liderando la carrera hacia la AGI", top=Inches(1.2), font_size=Pt(20), color=ACCENT_CYAN)

# Company info
bullets_company = [
    "Fundada en 2015 por Sam Altman, Elon Musk, Ilya Sutskever, Greg Brockman y otros como organización sin ánimo de lucro.",
    "Misión: 'Asegurar que la inteligencia artificial general beneficie a toda la humanidad'.",
    "2022: Lanzan ChatGPT, alcanzando 100 millones de usuarios en 2 meses — el crecimiento más rápido de la historia.",
    "GPT-4 (2023): modelo multimodal capaz de procesar texto e imágenes, superando el percentil 90 en exámenes como el BAR.",
    "Desarrollos: DALL·E (generación de imágenes), Codex (programación), Whisper (reconocimiento de voz), Sora (video).",
    "Estructura: OpenAI LP (limited profit) bajo OpenAI Inc. (non-profit). Microsoft es el principal inversor con $13B+.",
    "Valoración estimada: $157 mil millones (2025). Sede en San Francisco, California.",
]
add_body_text(slide, bullets_company, top=Inches(2.0), font_size=Pt(15.5))

# CEO bio
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.3), Inches(4.3), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
shape.line.color.rgb = ACCENT_NEON
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "👤 CEO: SAM ALTMAN"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_NEON
p.font.bold = True
p.font.name = 'Segoe UI'
p2 = tf.add_paragraph()
p2.text = "(n. 1985, Chicago). Emprendedor e inversor. Ex-presidente de Y Combinator (2014-2019). Co-fundador de Loopt y Worldcoin. Considerado una de las 100 personas más influyentes del mundo por Time (2023). Visionario que aboga por una AGI segura y regulada."
p2.font.size = Pt(11)
p2.font.color.rgb = LIGHT_GRAY
p2.font.name = 'Segoe UI Light'

add_page_number(slide, 15, TOTAL_SLIDES)

# ============================================================
# SLIDE 16: GOOGLE DEEPMIND
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_side_banner(slide, color=ACCENT_CYAN)
add_title_text(slide, "GOOGLE DEEPMIND", color=ACCENT_CYAN, font_size=Pt(44))
add_subtitle_text(slide, "Resolviendo la inteligencia para avanzar la ciencia", top=Inches(1.2), font_size=Pt(20), color=ACCENT_GOLD)

bullets = [
    "Fundada en 2010 en Londres por Demis Hassabis, Shane Legg y Mustafa Suleyman. Adquirida por Google en 2014 por $500+ millones.",
    "AlphaGo (2016): derrotó a Lee Sedol en Go. AlphaZero (2017): aprendió ajedrez, Go y shogi desde cero en horas, superando a todos los programas previos.",
    "AlphaFold (2020): resolvió el problema del plegamiento de proteínas, un desafío de 50 años en biología. Reconocido como uno de los mayores avances científicos del siglo.",
    "Gemini (2023/2024): modelo multimodal nativo que compite directamente con GPT-4, integrado en productos Google (Bard/Gemini, Search, Workspace).",
    "Investigación en IA general, robótica, salud (diagnóstico de cáncer de mama y enfermedades oculares), cambio climático y fusión nuclear.",
    "Google AI integra DeepMind y Google Brain (2023) bajo una sola división, consolidando su apuesta por ser líder en IA.",
]
add_body_text(slide, bullets, top=Inches(2.0), font_size=Pt(15.5))

# CEO bio
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.3), Inches(4.3), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
shape.line.color.rgb = ACCENT_CYAN
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "👤 CEO: DEMIS HASSABIS"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_CYAN
p.font.bold = True
p.font.name = 'Segoe UI'
p2 = tf.add_paragraph()
p2.text = "(n. 1976, Londres). Niño prodigio del ajedrez (Maestro a los 13). Neurocientífico por UCL y PhD en neurociencia cognitiva. Diseñador de videojuegos (Theme Park, Black & White). Fundó DeepMind con la misión de 'resolver la inteligencia'. Nombrado Sir en 2024 por servicios a la IA. Premio Nobel de Química 2024."
p2.font.size = Pt(11)
p2.font.color.rgb = LIGHT_GRAY
p2.font.name = 'Segoe UI Light'

add_page_number(slide, 16, TOTAL_SLIDES)

# ============================================================
# SLIDE 17: META AI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_side_banner(slide, color=ACCENT_CYAN)
add_title_text(slide, "META AI", color=ACCENT_CYAN, font_size=Pt(44))
add_subtitle_text(slide, "Open source y el metaverso impulsado por IA", top=Inches(1.2), font_size=Pt(20), color=ACCENT_GOLD)

bullets = [
    "Meta (anteriormente Facebook) creó FAIR (Facebook AI Research) en 2013, liderado por Yann LeCun, uno de los 'padrinos' del deep learning.",
    "Llama (Large Language Model Meta AI): familia de modelos open source. Llama 3 (2024) compite con GPT-4 y Gemini, disponible gratuitamente para investigación.",
    "Filosofía open source: Meta apuesta por democratizar la IA liberando modelos, pesos y herramientas. Esto acelera la innovación global.",
    "Investigación puntera en: NLP, visión artificial, robótica, traducción automática (No Language Left Behind), generación de código y vídeo.",
    "Integración en productos: Facebook, Instagram (filtros IA, recomendaciones), WhatsApp (asistente Meta AI), Ray-Ban Meta smart glasses.",
    "Inversión masiva en infraestructura: clusters con cientos de miles de GPUs NVIDIA H100 para entrenar la próxima generación de modelos.",
]
add_body_text(slide, bullets, top=Inches(2.0), font_size=Pt(15.5))

# CEO bio
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.3), Inches(4.3), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
shape.line.color.rgb = ACCENT_CYAN
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "👤 CEO: MARK ZUCKERBERG"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_CYAN
p.font.bold = True
p.font.name = 'Segoe UI'
p2 = tf.add_paragraph()
p2.text = "(n. 1984, Nueva York). Fundó Facebook en 2004 en Harvard. Visionario que transformó Meta hacia la IA y el metaverso. Defensor del open source en IA. Patrimonio neto: $200B+. Lidera personalmente la estrategia de IA de Meta, integrando Llama en todas las plataformas de la compañía."
p2.font.size = Pt(11)
p2.font.color.rgb = LIGHT_GRAY
p2.font.name = 'Segoe UI Light'

add_page_number(slide, 17, TOTAL_SLIDES)

# ============================================================
# SLIDE 18: MICROSOFT AI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_side_banner(slide, color=ACCENT_CYAN)
add_title_text(slide, "MICROSOFT AI", color=ACCENT_CYAN, font_size=Pt(44))
add_subtitle_text(slide, "El gigante que apostó por OpenAI y la IA integrada", top=Inches(1.2), font_size=Pt(20), color=ACCENT_GOLD)

bullets = [
    "Microsoft ha invertido más de $13 mil millones en OpenAI, asegurando acceso exclusivo a GPT-4 y modelos futuros para sus productos.",
    "Copilot: asistente de IA integrado en Windows, Office 365 (Word, Excel, PowerPoint, Teams), GitHub (Copilot para código) y Edge.",
    "Azure AI: plataforma cloud para entrenar y desplegar modelos de IA a escala empresarial. Infraestructura con GPUs NVIDIA y chips propios (Maia).",
    "Estrategia 'AI-first': 'Copilot en todo'. Cada producto Microsoft está siendo rediseñado con IA generativa integrada.",
    "Phi: familia de modelos pequeños y eficientes (SLMs) desarrollados por Microsoft Research que compiten en rendimiento con modelos mucho mayores.",
    "Investigación en: IA responsable, razonamiento, agentes autónomos, salud con IA, sostenibilidad y computación cuántica.",
]
add_body_text(slide, bullets, top=Inches(2.0), font_size=Pt(16))

# CEO bio
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.3), Inches(4.3), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
shape.line.color.rgb = ACCENT_CYAN
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "👤 CEO: SATYA NADELLA"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_CYAN
p.font.bold = True
p.font.name = 'Segoe UI'
p2 = tf.add_paragraph()
p2.text = "(n. 1967, Hyderabad, India). Ingeniero eléctrico y MBA. Tercer CEO de Microsoft (desde 2014). Transformó la empresa hacia cloud (Azure) e IA. Visionario que apostó temprano por OpenAI. Bajo su liderazgo, Microsoft superó los $3 billones de capitalización bursátil."
p2.font.size = Pt(11)
p2.font.color.rgb = LIGHT_GRAY
p2.font.name = 'Segoe UI Light'

add_page_number(slide, 18, TOTAL_SLIDES)

# ============================================================
# SLIDE 19: ANTHROPIC
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_side_banner(slide, color=ACCENT_CORAL)
add_title_text(slide, "ANTHROPIC", color=ACCENT_CORAL, font_size=Pt(44))
add_subtitle_text(slide, "Seguridad en la IA como prioridad fundamental", top=Inches(1.2), font_size=Pt(20), color=ACCENT_CYAN)

bullets = [
    "Fundada en 2021 por Dario Amodei y Daniela Amodei, ex-ejecutivos de OpenAI que abandonaron por diferencias en el enfoque de seguridad.",
    "Claude: familia de modelos de lenguaje (Claude 3.5 Sonnet, Opus, Haiku). Enfocados en ser 'helpful, honest, and harmless'.",
    "Constitutional AI: método propio de entrenamiento donde el modelo se alinea con una 'constitución' de principios éticos, no solo con preferencias humanas.",
    "Enfoque en seguridad: investigación puntera en interpretabilidad mecanicista (entender cómo 'piensan' los modelos) y alineamiento de la IA.",
    "Inversores: Google ($2B+), Amazon ($4B), Salesforce, Zoom. Valoración estimada: $60+ mil millones (2025).",
    "Considerados el principal competidor de OpenAI, con un enfoque diferencial en la seguridad y la investigación fundamental.",
]
add_body_text(slide, bullets, top=Inches(2.0), font_size=Pt(16))

# CEO bio
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.3), Inches(4.3), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
shape.line.color.rgb = ACCENT_CORAL
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "👤 CEO: DARIO AMODEI"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_CORAL
p.font.bold = True
p.font.name = 'Segoe UI'
p2 = tf.add_paragraph()
p2.text = "(n. 1983). PhD en Física por Princeton. Investigador postdoc en Stanford. Ex-VP de Investigación en OpenAI (2016-2020). Neurocientífico computacional de formación. Considerado una de las voces más autorizadas en seguridad de IA. Cree que la AGI llegará en esta década."
p2.font.size = Pt(11)
p2.font.color.rgb = LIGHT_GRAY
p2.font.name = 'Segoe UI Light'

add_page_number(slide, 19, TOTAL_SLIDES)

# ============================================================
# SLIDE 20: NVIDIA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_side_banner(slide, color=ACCENT_NEON)
add_title_text(slide, "NVIDIA", color=ACCENT_NEON, font_size=Pt(44))
add_subtitle_text(slide, "El hardware que hace posible la revolución de la IA", top=Inches(1.2), font_size=Pt(20), color=ACCENT_GOLD)

bullets = [
    "NVIDIA ha pasado de ser una empresa de videojuegos a convertirse en la empresa más valiosa del mundo ($3T+) gracias a sus GPUs para IA.",
    "Las GPUs NVIDIA (A100, H100, B200 Blackwell) son el estándar de facto para entrenar modelos de IA. El 80%+ del mercado de chips de IA.",
    "CUDA: plataforma de computación paralela que permite a los desarrolladores programar GPUs. El 'foso competitivo' de NVIDIA con +4M de desarrolladores.",
    "Todos los grandes modelos (GPT-4, Gemini, Llama, Claude) se entrenan en clusters con decenas de miles de GPUs NVIDIA H100.",
    "DGX SuperPOD: supercomputadores de IA llave en mano. NVIDIA vende infraestructura completa (hardware + software + networking).",
    "Inversión en robótica (Jetson, Isaac), coches autónomos (DRIVE), salud (Clara), omniverse (gemelos digitales) y computación cuántica.",
]
add_body_text(slide, bullets, top=Inches(2.0), font_size=Pt(15.5))

# CEO bio
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.3), Inches(4.3), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
shape.line.color.rgb = ACCENT_NEON
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "👤 CEO: JENSEN HUANG"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_NEON
p.font.bold = True
p.font.name = 'Segoe UI'
p2 = tf.add_paragraph()
p2.text = "(n. 1963, Taiwán). Ingeniero eléctrico por Stanford. Co-fundó NVIDIA en 1993. Visionario que apostó por las GPUs para IA en 2006 con CUDA. Ícono cultural (chaqueta de cuero negra). Bajo su liderazgo, NVIDIA creció de $10B a $3T. Una de las personas más ricas del mundo."
p2.font.size = Pt(11)
p2.font.color.rgb = LIGHT_GRAY
p2.font.name = 'Segoe UI Light'

add_page_number(slide, 20, TOTAL_SLIDES)

# ============================================================
# SLIDE 21: xAI / TESLA AI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_side_banner(slide, color=ACCENT_GOLD)
add_title_text(slide, "TESLA AI & xAI", color=ACCENT_GOLD, font_size=Pt(44))
add_subtitle_text(slide, "El visionario controvertido: coches autónomos y búsqueda de la verdad", top=Inches(1.2), font_size=Pt(20), color=ACCENT_CYAN)

bullets = [
    "Tesla AI: desarrolla el sistema Full Self-Driving (FSD) con redes neuronales entrenadas con datos de millones de vehículos reales.",
    "Dojo: supercomputador propio de Tesla para entrenamiento de IA de conducción autónoma usando chips diseñados internamente.",
    "Optimus: robot humanoide de Tesla impulsado por la misma IA de FSD. Diseñado para tareas repetitivas en fábricas y hogares.",
    "xAI: fundada por Elon Musk en 2023. Misión: 'entender la verdadera naturaleza del universo'. Desarrolla Grok, asistente integrado en X (Twitter).",
    "Grok: modelo de IA con 'personalidad' y sentido del humor. Acceso en tiempo real a X. Enfoque en razonamiento y búsqueda de la verdad.",
    "Musk fue co-fundador de OpenAI (2015) y la abandonó en 2018. Crítico de la dirección actual de OpenAI, aboga por IA de código abierto.",
]
add_body_text(slide, bullets, top=Inches(2.0), font_size=Pt(15))

# CEO bio
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.3), Inches(4.3), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
shape.line.color.rgb = ACCENT_GOLD
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "👤 CEO: ELON MUSK"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_GOLD
p.font.bold = True
p.font.name = 'Segoe UI'
p2 = tf.add_paragraph()
p2.text = "(n. 1971, Pretoria, Sudáfrica). Fundador de Tesla, SpaceX, xAI, Neuralink y The Boring Company. Co-fundó OpenAI en 2015. El hombre más rico del mundo ($300B+). Visionario polarizante pero innegablemente influyente en la IA moderna."
p2.font.size = Pt(11)
p2.font.color.rgb = LIGHT_GRAY
p2.font.name = 'Segoe UI Light'

add_page_number(slide, 21, TOTAL_SLIDES)

# ============================================================
# SLIDE 22: APPLE AI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_side_banner(slide, color=ACCENT_CYAN)
add_title_text(slide, "APPLE INTELLIGENCE", color=ACCENT_CYAN, font_size=Pt(44))
add_subtitle_text(slide, "Privacidad y la IA en el ecosistema Apple", top=Inches(1.2), font_size=Pt(20), color=ACCENT_GOLD)

bullets = [
    "Apple Intelligence (anunciada en WWDC 2024): integración de IA generativa en iPhone, iPad y Mac con procesamiento on-device como prioridad.",
    "Private Cloud Compute: servidores Apple Silicon que extienden la capacidad de IA garantizando privacidad. Los datos no se almacenan ni se usan para entrenar.",
    "Siri renovada: comprensión contextual avanzada, acciones entre apps, conciencia en pantalla y capacidad de delegar a ChatGPT.",
    "Modelos fundacionales propios: Apple desarrolla modelos de IA optimizados para ejecutarse en dispositivos móviles con Neural Engine.",
    "Filosofía: la IA debe ser personal, privada y útil. Procesamiento en el dispositivo como diferencial competitivo frente a Google y Microsoft.",
    "Investigación en: ML eficiente, visión artificial, NLP on-device, salud digital (Apple Watch), realidad aumentada y Spatial Computing.",
]
add_body_text(slide, bullets, top=Inches(2.0), font_size=Pt(16))

# CEO bio
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.3), Inches(4.3), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
shape.line.color.rgb = ACCENT_CYAN
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "👤 CEO: TIM COOK"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_CYAN
p.font.bold = True
p.font.name = 'Segoe UI'
p2 = tf.add_paragraph()
p2.text = "(n. 1960, Alabama, EE.UU.). MBA por Duke. CEO de Apple desde 2011, sucediendo a Steve Jobs. Experto en cadena de suministro. Bajo su liderazgo, Apple alcanzó $3T de valoración. Enfoque en privacidad como derecho humano fundamental."
p2.font.size = Pt(11)
p2.font.color.rgb = LIGHT_GRAY
p2.font.name = 'Segoe UI Light'

add_page_number(slide, 22, TOTAL_SLIDES)

# ============================================================
# SLIDE 23: AMAZON AI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_side_banner(slide, color=ACCENT_CORAL)
add_title_text(slide, "AMAZON AI", color=ACCENT_CORAL, font_size=Pt(44))
add_subtitle_text(slide, "Alexa, AWS y la IA a escala planetaria", top=Inches(1.2), font_size=Pt(20), color=ACCENT_CYAN)

bullets = [
    "AWS (Amazon Web Services): la plataforma cloud líder mundial ofrece servicios de IA/ML (SageMaker, Bedrock, Rekognition, Polly, Lex).",
    "Alexa: asistente de voz presente en cientos de millones de dispositivos. En renovación con Alexa LLM para conversaciones naturales y proactivas.",
    "Amazon Bedrock: servicio serverless para acceder a modelos fundacionales de Anthropic, Meta, Mistral, Stability AI y los propios de Amazon (Titan).",
    "Inversión estratégica: $4 mil millones en Anthropic, desarrollo de chips propios (Trainium para entrenamiento, Inferentia para inferencia).",
    "Aplicaciones internas: logística predictiva, recomendaciones (35% de ventas), centros de distribución robotizados, drones Prime Air, Amazon Go.",
    "Proyecto Kuiper: constelación de satélites para internet global, integrará IA para optimizar comunicaciones y cobertura.",
]
add_body_text(slide, bullets, top=Inches(2.0), font_size=Pt(15.5))

# CEO bio
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.3), Inches(4.3), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
shape.line.color.rgb = ACCENT_CORAL
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "👤 CEO: ANDY JASSY"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_CORAL
p.font.bold = True
p.font.name = 'Segoe UI'
p2 = tf.add_paragraph()
p2.text = "(n. 1968, Nueva York). MBA por Harvard. CEO de Amazon desde 2021, sucediendo a Jeff Bezos. Arquitecto de AWS, la división que genera la mayor parte de beneficios de Amazon. Enfoque en IA generativa aplicada a escala empresarial."
p2.font.size = Pt(11)
p2.font.color.rgb = LIGHT_GRAY
p2.font.name = 'Segoe UI Light'

add_page_number(slide, 23, TOTAL_SLIDES)

# ============================================================
# SLIDE 24: OTRAS EMPRESAS CLAVE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_title_text(slide, "OTRAS EMPRESAS CLAVE EN IA", color=ACCENT_PURPLE, font_size=Pt(38))

companies = [
    ("Hugging Face", ACCENT_GOLD, "Clement Delangue", "La 'GitHub de la IA'. Plataforma colaborativa con cientos de miles de modelos open source. Hub central de la comunidad de ML."),
    ("Mistral AI", ACCENT_CYAN, "Arthur Mensch", "Startup francesa que compite con modelos open source de alto rendimiento. Mistral Large compite con GPT-4. Valoración: $6B+."),
    ("Stability AI", ACCENT_PURPLE, "Prem Akkaraju", "Creadores de Stable Diffusion, el modelo de generación de imágenes open source más popular del mundo."),
    ("Cohere", ACCENT_CORAL, "Aidan Gomez", "Enfocada en IA para empresas. Modelos de lenguaje optimizados para RAG, búsqueda y embeddings. Fundada por autores del paper Transformer."),
    ("Scale AI", ACCENT_NEON, "Alexandr Wang", "Plataforma de datos para entrenar IA. Provee datos etiquetados de alta calidad a OpenAI, Google, Meta y el Departamento de Defensa de EE.UU."),
    ("Perplexity AI", ACCENT_GOLD, "Aravind Srinivas", "Buscador conversacional con IA que desafía a Google. Combina modelos de lenguaje con búsqueda en tiempo real. Valoración: $9B+."),
]

y = Inches(1.7)
for name, color, ceo, desc in companies:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(2.6), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = f"{name} — {ceo}"
    p.font.size = Pt(12)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER
    # Description
    txBox = slide.shapes.add_textbox(Inches(3.6), y, Inches(9), Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = 'Segoe UI Light'
    y += Inches(0.85)

add_page_number(slide, 24, TOTAL_SLIDES)

# ============================================================
# SLIDE 25: LÍNEA DE TIEMPO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_title_text(slide, "LÍNEA DE TIEMPO DE LA IA", color=ACCENT_GOLD, font_size=Pt(40))

timeline = [
    ("1950", "Test de Turing", ACCENT_CYAN),
    ("1956", "Conferencia Dartmouth", ACCENT_GOLD),
    ("1958", "Lenguaje LISP", ACCENT_PURPLE),
    ("1966", "Chatbot ELIZA", ACCENT_CORAL),
    ("1973", "Invierno IA", ACCENT_NEON),
    ("1980", "Sistemas Expertos", ACCENT_CYAN),
    ("1986", "Backpropagation", ACCENT_GOLD),
    ("1997", "Deep Blue vence a Kasparov", ACCENT_PURPLE),
    ("2006", "Deep Learning (Hinton)", ACCENT_CORAL),
    ("2011", "Siri / Watson", ACCENT_NEON),
    ("2012", "AlexNet (Big Bang DL)", ACCENT_CYAN),
    ("2014", "GANs / DeepMind", ACCENT_GOLD),
    ("2016", "AlphaGo vence a Lee Sedol", ACCENT_PURPLE),
    ("2017", "Transformers", ACCENT_CORAL),
    ("2020", "GPT-3 / AlphaFold", ACCENT_NEON),
    ("2022", "ChatGPT (100M usuarios)", ACCENT_CYAN),
    ("2024", "GPT-4o, Gemini, Claude 3.5", ACCENT_GOLD),
    ("2025+", "¿AGI?", ACCENT_GOLD),
]

# Draw timeline line
add_banner(slide, Inches(1.2), Inches(3.7), Inches(11), Inches(0.03), ACCENT_GOLD)

# Place nodes on zigzag
for i, (year, event, color) in enumerate(timeline):
    if i < 9:
        y = Inches(2.8)
    else:
        y = Inches(4.6)
    x = Inches(0.8) + Inches(i % 9) * Inches(1.35)
    if i >= 9:
        x = Inches(0.8) + Inches((i - 9)) * Inches(1.35)

    # Dot
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), Inches(3.6), Inches(0.15), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    # Year
    txBox = slide.shapes.add_textbox(x - Inches(0.1), y - Inches(0.1), Inches(0.9), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = year
    p.font.size = Pt(11)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER

    # Event
    txBox = slide.shapes.add_textbox(x - Inches(0.3), y + Inches(0.25), Inches(1.45), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = event
    p.font.size = Pt(9)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = 'Segoe UI Light'
    p.alignment = PP_ALIGN.CENTER

add_page_number(slide, 25, TOTAL_SLIDES)

# ============================================================
# SLIDE 26: EL FUTURO DE LA IA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_title_text(slide, "EL FUTURO DE LA IA", color=ACCENT_NEON, font_size=Pt(44))

bullets = [
    "🤖  AGI (Inteligencia Artificial General): predicciones de líderes como Sam Altman y Demis Hassabis apuntan a esta década. Un sistema capaz de realizar cualquier tarea intelectual humana.",
    "🧠  Neurociencia + IA: interfaces cerebro-computadora (Neuralink), comprensión de la conciencia y modelado del cerebro humano.",
    "🏥  Revolución en salud: diagnóstico temprano, descubrimiento de fármacos asistido por IA, medicina personalizada y longevidad extendida.",
    "🌍  Cambio climático: optimización de redes energéticas, captura de carbono, fusión nuclear acelerada por IA y predicción de desastres naturales.",
    "🎓  Educación personalizada: tutores de IA adaptados al ritmo y estilo de aprendizaje de cada estudiante, democratizando la educación global.",
    "⚠️  Desafíos: desempleo tecnológico, desinformación generada por IA, sesgos algorítmicos, privacidad, concentración de poder y el problema del alineamiento.",
    "🔬  Descubrimiento científico: la IA como científico autónomo, generando hipótesis y diseñando experimentos. AlphaFold ya demostró este potencial.",
]
add_body_text(slide, bullets, top=Inches(1.8), font_size=Pt(18))

# Highlight
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.9))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x0A)
shape.line.color.rgb = ACCENT_NEON
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🌟  El mayor desafío: asegurar que la IA beneficie a toda la humanidad, no solo a unos pocos. La gobernanza de la IA será el tema definitorio del siglo XXI."
p.font.size = Pt(17)
p.font.color.rgb = ACCENT_NEON
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

add_page_number(slide, 26, TOTAL_SLIDES)

# ============================================================
# SLIDE 27: CONCLUSIONES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_gradient_banner(slide)
add_title_text(slide, "CONCLUSIONES", color=ACCENT_GOLD, font_size=Pt(44))

bullets = [
    "La IA no nació de la noche a la mañana: es el resultado de 70+ años de investigación, fracasos y avances de miles de científicos.",
    "Los precursores (Turing, McCarthy, Minsky, Shannon) sentaron las bases teóricas décadas antes de que la tecnología pudiera implementarlas.",
    "Hemos pasado de sistemas basados en reglas a redes neuronales que aprenden de datos, acercándonos cada vez más a una inteligencia flexible.",
    "Las grandes empresas actuales (OpenAI, Google DeepMind, Meta, Microsoft, Anthropic, NVIDIA) lideran una revolución sin precedentes.",
    "Los CEOs actuales no son solo ejecutivos: son visionarios que moldean el futuro de la humanidad. Sus decisiones afectarán a miles de millones.",
    "Estamos en un punto de inflexión histórico. La IA está transformando la ciencia, la economía, el arte y la sociedad a una velocidad exponencial.",
    "La pregunta ya no es si las máquinas pueden pensar, sino cómo aseguramos que piensen en beneficio de todos.",
]
add_body_text(slide, bullets, top=Inches(1.8), font_size=Pt(18))

# Quote
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
shape.line.color.rgb = ACCENT_GOLD
shape.line.width = Pt(1)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"La inteligencia artificial es probablemente lo más importante que la humanidad haya inventado jamás. Pienso que es más profundo que el fuego o la electricidad." — Sundar Pichai, CEO de Google'
p.font.size = Pt(15)
p.font.color.rgb = ACCENT_GOLD
p.font.italic = True
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

add_page_number(slide, 27, TOTAL_SLIDES)

# ============================================================
# SLIDE 28: GRACIAS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

# Top banner
add_banner(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), ACCENT_CYAN)
add_banner(slide, Inches(0), Inches(0.12), Inches(13.333), Inches(0.04), ACCENT_GOLD)
add_banner(slide, Inches(0), Inches(0.16), Inches(13.333), Inches(0.04), ACCENT_CORAL)
add_banner(slide, Inches(0), Inches(0.20), Inches(13.333), Inches(0.04), ACCENT_PURPLE)

add_title_text(slide, "GRACIAS", top=Inches(2.0), font_size=Pt(64), color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER, width=Inches(12))
add_subtitle_text(slide, "Historia de la Inteligencia Artificial", top=Inches(3.2), font_size=Pt(26), color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, width=Inches(12))
add_subtitle_text(slide, "De los precursores a las grandes empresas del siglo XXI", top=Inches(3.8), font_size=Pt(18), color=MED_GRAY, alignment=PP_ALIGN.CENTER, width=Inches(12))

add_banner(slide, Inches(4), Inches(4.5), Inches(5.333), Inches(0.03), ACCENT_GOLD)

add_subtitle_text(slide, "Presentación generada por", top=Inches(5.0), font_size=Pt(17), color=MED_GRAY, alignment=PP_ALIGN.CENTER, width=Inches(12))
add_title_text(slide, "JARVIS · Agente NEO", top=Inches(5.5), font_size=Pt(30), color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER, width=Inches(12))
add_subtitle_text(slide, "Inteligencia Artificial al servicio de la humanidad · 2026", top=Inches(6.3), font_size=Pt(15), color=DARK_GRAY, alignment=PP_ALIGN.CENTER, width=Inches(12))

# Bottom banner
add_banner(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.04), ACCENT_PURPLE)
add_banner(slide, Inches(0), Inches(7.39), Inches(13.333), Inches(0.04), ACCENT_CORAL)
add_banner(slide, Inches(0), Inches(7.43), Inches(13.333), Inches(0.04), ACCENT_GOLD)
add_banner(slide, Inches(0), Inches(7.47), Inches(13.333), Inches(0.08), ACCENT_CYAN)

add_page_number(slide, 28, TOTAL_SLIDES)

# ============================================================
# SAVE
# ============================================================
output_dir = "/home/dorti/jarvis-voice-factory/projects"
output_file = os.path.join(output_dir, "historia_de_la_inteligencia_artificial.pptx")
os.makedirs(output_dir, exist_ok=True)
prs.save(output_file)
print(f"✅ Presentación guardada: {output_file}")
print(f"📊 Total slides: {len(prs.slides)}")
